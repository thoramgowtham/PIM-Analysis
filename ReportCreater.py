import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import sys
from openpyxl import Workbook
from openpyxl import load_workbook
import traceback
from tkinter import *
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

#Execution : python Final_Part1_V1.py FALOCATION SITEID Sectors

left = Inches(0.69)
top = Inches(2.07)
width = Inches(4.64)
height = Inches(4.64)
sectors = sys.argv[3]


#Definitions for table borders
def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def _set_cell_border(cell, border_color="000000", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
#definition to get slide count
def getSlidecount():
	
	x = 0
	slides = prs.slides
	for slide in slides:
		x = x + 1
	return x
	
		
#function to get c7 and c4 comments this should return direct comment with score specific
def GetScoringComment(picnamewithoutExt,cat4or7):
	picnamewithoutExtSplit = (picnamewithoutExt.split("_"))
	PicSector = str(picnamewithoutExtSplit[1])[0]
	if PicSector == 'G':
		PicSector = 'C'
	try:
		for i in range(2,sectors + 2):
			slideTitle = prs.slides[i].shapes[2].text_frame.paragraphs[0].text
			print(slideTitle)
			slideTitle = slideTitle.split("_")
			print(slideTitle)
			slideSector = str(slideTitle[2])[1]
			print('slide sector : ' + ' ' + slideSector +'  ' + 'picSector : '+ ' '+ PicSector)
			if (str(slideSector.upper()) == str(PicSector.upper())):
				print(slideSector)
				SlideTableGSC = prs.slides[i].shapes[7].table
				C7Score = str(SlideTableGSC.cell(1,6).text_frame.paragraphs[0].text)
				C4Score = str(SlideTableGSC.cell(1,3).text_frame.paragraphs[0].text)
				if cat4or7 == 'c7':
					if C7Score == '1':
						return '• Skewed antenna > 15°'
					elif C7Score == '2':
						return '• Skewed antenna > 30°'
					elif C7Score == '3':
						return '• Skewed antenna > 45°'
					else:
						return '• Skewed antenna > 15°'
				if cat4or7 == 'c4':
					if C4Score == '1':
						return '• Antenna separation <= 6ft'
					elif C4Score == '2':
						return '• Antenna separation <= 3ft'
					elif C4Score == '3':
						return '• Antenna separation <= 1ft'
					else:
						return '• Antenna separation < 6ft'
			else:
				print('C4 and C7 Scoring is not compatible with Excel Scoring given')
	except:
		return ' '
		
	
		
		
#Comments for the pictures
def GetComment(cat,picnamewithoutExt):
	picnamewithoutExt = str(picnamewithoutExt) 
	if cat.upper() == 'C1':
		return '• Contamination'
	if cat.upper() == 'C2':
		return '• Loose connectors'
	if cat.upper() == 'C3':
		return '• Equipment right behind the antenna/Bent wires'
	if cat.upper() == 'C4':
		return GetScoringComment(picnamewithoutExt,'c4')
	if cat.upper() == 'C5':
		return '• Blocking in front of Antenna'
	if cat.upper() == 'C6':
		return '• Hardware Broken/Faulty'
	if cat.upper() == 'C7':
		return GetScoringComment(picnamewithoutExt,'c7')
	if cat.upper() == 'C8':
		return '• Feedline cable clutter'
	if cat.upper() == 'C9':
		return '• Equipment need to be Tighten'
	if cat.upper() == 'C10':
		return '• Consider replacing/rearranging cables'
	if cat.upper() == 'C11':
		return '• Stacked Hangers'
	if cat.upper() == 'C12':
		return '• Miscellaneous'
	return

def GetPicComment(name,position): #position 1 if comment for 1st pic in a slide or 2 for second pic in a slide
	picname = name
	picnamewithoutExt =picname.split(".")
	picnamesplit = (picnamewithoutExt[0].split("_"))
	if position == '1': 
		txbox = Slide.shapes.add_textbox(Inches(0.46), Inches(6.83), Inches(5.19), Inches(0.52))
	if position == '2':
		txbox = Slide.shapes.add_textbox(Inches(6.4), Inches(6.83), Inches(5.19), Inches(0.52))
	if len(picnamesplit) <= 2:
		txbox.text_frame.paragraphs[0].text = '• No issues Found'
		txbox.text_frame.paragraphs[0].font.size = Pt(14)
	else:
		txbox.text_frame.paragraphs[0].text = GetComment(str(picnamesplit[2]),str(picnamewithoutExt[0])) #str(picnamesplit[2]) is 'C1' to 'C12' string
		txbox.text_frame.paragraphs[0].font.size = Pt(14)
	return

#Functions for Site visits
def getSiteVisitSlides(sec):
	print('getSiteVisitSlides Called')
	blank_slide_layout = prs.slide_layouts[6]
	Slide = prs.slides.add_slide(blank_slide_layout)
	pic = Slide.shapes.add_picture(att_path,Inches(12.31),Inches(6.43),Inches(0.92),Inches(0.92))
	pic = Slide.shapes.add_picture(tts_path, Inches(0),Inches(0),Inches(13.33),Inches(0.92))
	
	if sec == 'A':
		sector = 'Alpha'
	elif sec == 'B':
		sector = 'Beta'
	elif sec == 'C':
		sector = 'Gamma'
	
	
	txBox = Slide.shapes.add_textbox(Inches(0.25), Inches(0.65), Inches(8), Inches(0.7))
	txBox.text_frame.paragraphs[0].text = 'Spectrum Sweeps – ' +  sector  #here sector is 'Alpha' or 'Beta' or 'Gamma'
	txBox.text_frame.paragraphs[0].font.size = Pt(32)

	
	sector = str(sec)
	
	try:
		pic1_path = sitevisitpath + faloc +  ' FirstNet Site Sweeps '+ sector +'/' + faloc + '.' + sector + '.D.B.1314.jpg'
		pic2_path = sitevisitpath + faloc +  ' FirstNet Site Sweeps '+ sector +'/' + faloc + '.' + sector + '.U.B.1314b.jpg'
		pic3_path = sitevisitpath + faloc +  ' FirstNet Site Sweeps '+ sector +'/' + faloc + '.' + sector + '.D.B.2912.jpg'
		pic4_path = sitevisitpath + faloc +  ' FirstNet Site Sweeps '+ sector +'/' + faloc + '.' + sector + '.D.B.5.jpg'
		pic5_path = sitevisitpath + faloc +  ' FirstNet Site Sweeps '+ sector +'/' + faloc + '.' + sector + '.D.B.25.jpg'
		pic6_path = sitevisitpath + faloc +  ' FirstNet Site Sweeps '+ sector +'/' + faloc + '.' + sector + '.D.B.66.jpg'
		pic7_path = sitevisitpath + faloc +  ' FirstNet Site Sweeps '+ sector +'/' + faloc + '.' + sector + '.D.B.30.jpg'
		
		#---- Please add text box posisions after consulting tg team ! ----
		txBox = Slide.shapes.add_textbox(Inches(1.83), Inches(3.72), Inches(1.12), Inches(0.4))
		txBox.text_frame.paragraphs[0].text = 'D.B1314'
		txBox.text_frame.paragraphs[0].font.size = Pt(18)
		txBox = Slide.shapes.add_textbox(Inches(6.04), Inches(3.72), Inches(1.12), Inches(0.4))
		txBox.text_frame.paragraphs[0].text = 'U.B1314'
		txBox.text_frame.paragraphs[0].font.size = Pt(18)
		txBox = Slide.shapes.add_textbox(Inches(10.25), Inches(3.72), Inches(1.12), Inches(0.4))
		txBox.text_frame.paragraphs[0].text = 'D.B2912'
		txBox.text_frame.paragraphs[0].font.size = Pt(18)
		
		slideCount = getSlidecount() - 1
		
		#sector 1st slide Pictures
		try:
			pic = Slide.shapes.add_picture(pic1_path,Inches(0.67),Inches(1.39),Inches(3.63),Inches(2.29))
			pic = Slide.shapes.add_picture(pic2_path,Inches(4.89),Inches(1.39),Inches(3.63),Inches(2.29))
			pic = Slide.shapes.add_picture(pic3_path,Inches(9.1),Inches(1.39),Inches(3.63),Inches(2.29))
			
			getMarkers(slideCount)
			
			
			#---- Please add text box posisions after consulting tg team ! ----
			
			if sector == 'A':
				sector = 'Alpha'
			elif sector == 'B':
				sector = 'Beta'
			elif sector == 'C':
				sector = 'Gamma'
			
			#sector 2nd slide
			blank_slide_layout = prs.slide_layouts[6]
			Slide = prs.slides.add_slide(blank_slide_layout)
			pic = Slide.shapes.add_picture(att_path,Inches(12.31),Inches(6.43),Inches(0.92),Inches(0.92))
			pic = Slide.shapes.add_picture(tts_path, Inches(0),Inches(0),Inches(13.33),Inches(0.92))
			txBox = Slide.shapes.add_textbox(Inches(0.25), Inches(0.65), Inches(8), Inches(0.7))
			txBox.text_frame.paragraphs[0].text = 'Spectrum Sweeps – ' +  sector  #here sec is 'Alpha' or 'Beta' or 'Gamma'
			txBox.text_frame.paragraphs[0].font.size = Pt(32)
			
			pic = Slide.shapes.add_picture(pic4_path,Inches(1.47),Inches(1.33),Inches(3.93),Inches(2.29))
			pic = Slide.shapes.add_picture(pic5_path,Inches(7.05),Inches(1.33),Inches(3.93),Inches(2.29))
			pic = Slide.shapes.add_picture(pic6_path,Inches(1.47),Inches(4.56),Inches(3.93),Inches(2.29))
			pic = Slide.shapes.add_picture(pic7_path,Inches(7.05),Inches(4.56),Inches(3.93),Inches(2.29))	
			#---- Please add text box posisions after consulting tg team ! ----
			txBox = Slide.shapes.add_textbox(Inches(2.71), Inches(3.68), Inches(1.12), Inches(0.4))
			txBox.text_frame.paragraphs[0].text = 'D.B5'
			txBox.text_frame.paragraphs[0].font.size = Pt(18)
			txBox = Slide.shapes.add_textbox(Inches(8.25), Inches(3.68), Inches(1.12), Inches(0.4))
			txBox.text_frame.paragraphs[0].text = 'D.B25'
			txBox.text_frame.paragraphs[0].font.size = Pt(18)
			txBox = Slide.shapes.add_textbox(Inches(2.71), Inches(6.97), Inches(1.12), Inches(0.4))
			txBox.text_frame.paragraphs[0].text = 'D.B66'
			txBox.text_frame.paragraphs[0].font.size = Pt(18)
			txBox = Slide.shapes.add_textbox(Inches(8.25), Inches(6.97), Inches(1.12), Inches(0.4))
			txBox.text_frame.paragraphs[0].text = 'D.B30'
			txBox.text_frame.paragraphs[0].font.size = Pt(18)
				
		except:
			pass	
			traceback.print_exc()
		#---- Chain reaction --> adding comments to the current slide ---- ! 
		
		getSpectrumComments(slideCount,sec)
		
		
		#---- Please add text box posisions after consulting tg team ! ----
		
	except Exception as e:
		ws1.cell(row = ws1.max_row + 1 , column = 1).value = sys.argv[1]
		ws1.cell(row = ws1.max_row  , column = 2).value = 'No pictures found'
		ws1.cell(row = ws1.max_row  , column = 3).value = str(e)
		traceback.print_exc()
		pass
	
	
	print('end of funtion call')

#funtion to get site visit spectrum pics comments
def getSpectrumComments(slide,sector):
	print('getSpectrumComments function called')
	#comment box 1
	slide = prs.slides[slide]
	if sector == 'A':
		CommentB1314 = ACommentB1314
		CommentB1314b = ACommentB1314b
		CommentB2912 = ACommentB2912
	if sector == 'B':
		CommentB1314 = BCommentB1314
		CommentB1314b = BCommentB1314b
		CommentB2912 = BCommentB2912
		
	if sector == 'C':
		CommentB1314 = CCommentB1314
		CommentB1314b = CCommentB1314b
		CommentB2912 = CCommentB2912
	
	txtbox1 = slide.shapes.add_textbox(Inches(0.51),Inches(4.12),Inches(3.98),Inches(3.23))
	txtbox1.line.color.rgb = RGBColor(0,0,0)
	txtbox1.text_frame.margin_bottom = Inches(0.10)
	txtbox1.text_frame.word_wrap = True
	
	txtbox2 = slide.shapes.add_textbox(Inches(4.68),Inches(4.12),Inches(3.98),Inches(3.23))
	txtbox2.line.color.rgb = RGBColor(0,0,0)
	txtbox2.text_frame.margin_bottom = Inches(0.10)
	txtbox2.text_frame.word_wrap = True
	
	txtbox3 = slide.shapes.add_textbox(Inches(8.82),Inches(4.12),Inches(3.98),Inches(3.23))
	txtbox3.line.color.rgb = RGBColor(0,0,0)
	txtbox3.text_frame.margin_bottom = Inches(0.10)
	txtbox3.text_frame.word_wrap = True
	
	# Comment for picture 1
	if CommentB1314 == 'VZW Colo - As expected':
		print('Case 1 identified ')
		txtbox1.text_frame.paragraphs[0].text = u'\u2022 Verizon Owned B13 can be additional source of PIM products'
		txtbox1.text_frame.add_paragraph()
		txtbox1.text_frame.paragraphs[1].text = u'\u2022 Sites co-located with Verizon needs extra attention'
		txtbox1.text_frame.add_paragraph()
		txtbox1.text_frame.paragraphs[2].text =  u'\u2022 B13 mixing can result in more severe IM3 & IM5'
	elif CommentB1314 == 'No Verizon colocation':
		txtbox1.text_frame.paragraphs[0].text = u'\u2022 No Verizon Colocation '
	elif CommentB1314 == 'VZW Colo - Further investigation':
		txtbox1.text_frame.paragraphs[0].text = u'\u2022 Verizon Co-located and spill over observed into B14'
		txtbox1.text_frame.add_paragraph()
	else:
		txtbox1.text_frame.paragraphs[0].text = u'\u2022' + CommentB1314
		
		
	#Comment picture 2
	if CommentB1314b == 'As expected':
		txtbox2.text_frame.paragraphs[0].text = u'\u2022 As Expected'
	elif CommentB1314b == 'Inconclusive':
		txtbox2.text_frame.paragraphs[0].text = u'\u2022 Inconclusive'
	else:
		txtbox2.text_frame.paragraphs[0].text = u'\u2022' + CommentB1314b
		
	#Comment for picture 3
	if CommentB2912 == 'RSSI Levels observed in Band 12':
		txtbox3.text_frame.paragraphs[0].text = u'\u2022 RSSI levels observed in Band 12'
		txtbox3.text_frame.add_paragraph()
		txtbox3.text_frame.paragraphs[1].text = u'\u2022 Band combinations of DL B17, B14 & B29 have potential to produce PIM in both the UL of B17 & B14'
		txtbox3.text_frame.add_paragraph()
		txtbox3.text_frame.paragraphs[2].text =  u'\u2022 Severity of IM3,IM5 & IM7'
		txtbox3.text_frame.add_paragraph()
		txtbox3.text_frame.paragraphs[3].text =  u'\u2022 Self PIM (DL Source within the same Band)'
	else:
		txtbox3.text_frame.paragraphs[0].text = u'\u2022' + CommentB2912
		
#Getting markers for the slide
def getMarkers(Slide):
	Slide = prs.slides[Slide]
	#------------- ON WORKS -------------------------------------------------------------------------------------------------------------------------
	line = Slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.77), Inches(1.66), Inches(0.005), Inches (1.42))
	line.line.fill.solid()	
	line.line.fill.fore_color.rgb = RGBColor(255,0,0)
	
	line = Slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.9), Inches(1.66), Inches(0.005), Inches (1.42))
	line.line.fill.solid()	
	line.line.fill.fore_color.rgb = RGBColor(255,0,0)
	
	line = Slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.18), Inches(1.66), Inches(0.005), Inches (1.36))
	line.line.fill.solid()	
	line.line.fill.fore_color.rgb = RGBColor(255,0,0)
	
	line = Slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.06), Inches(1.66), Inches(0.005), Inches (1.36))
	line.line.fill.solid()	
	line.line.fill.fore_color.rgb = RGBColor(255,0,0)
	
	
	line = Slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(11.11), Inches(1.66), Inches(0.005), Inches (1.36))
	line.line.fill.solid()	
	line.line.fill.fore_color.rgb = RGBColor(255,0,0)
	
	line = Slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(11.61), Inches(1.66), Inches(0.005), Inches (1.36))
	line.line.fill.solid()	
	line.line.fill.fore_color.rgb = RGBColor(255,0,0)
	
	line = Slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(12.11), Inches(1.66), Inches(0.005), Inches (1.36))
	line.line.fill.solid()	
	line.line.fill.fore_color.rgb = RGBColor(255,0,0)
	#------------------------------------------------------------------------------------------------------------------------------------------------
		
	

	
	
# .........Variables ...........	

market = 'PHASE 1 SITE VISIT'        #CHANGE MARKET
dir_path = 'C:/Users/me/Desktop/PPTCODES/Source/ATT-PIM-PHASE1-AGYLE2-MARKED/' + market + '/' + sys.argv[1]   
rfds_path = 'C:/Users/me/Desktop/PPTCODES/Source/ATT-PIM-PHASE1-AGYLE3-RFDS & SITEPLANS/' + market +'/' + sys.argv[1] 
att_path = 'C:/Users/me/Desktop/PPTCODES/Source/ATT-TTS-ICONS/att_icon.png'
tts_path = 'C:/Users/me/Desktop/PPTCODES/Source/ATT-TTS-ICONS/tts_icon.png'
Gepic1_path = 'C:/Users/me/Desktop/PPTCODES/Source/GOOGLE EARTH PICTURES/'+ sys.argv[1] +'/' +sys.argv[1] +'_GE.png'
Gepic2_path = 'C:/Users/me/Desktop/PPTCODES/Source/GOOGLE EARTH PICTURES/' + sys.argv[1] +'/' +sys.argv[1] +'_GESV.png'
defpath = 'C:/Users/me/Desktop/PPTCODES/Source/ATT-TTS-ICONS/UpdatedDef.png'
thankyoupicpath  = 'C:/Users/me/Desktop/PPTCODES/Source/ATT-TTS-ICONS/thankyou.png'
sitevisitpath = 'C:/Users/me/Desktop/PPTCODES/Source/ATT-PIM-PHASE1-SITEVISITS-SPECTRUM/' + sys.argv[1] + '/' + sys.argv[1]  + '/'


#debugger log file
debugger = 'C:/Users/me/Desktop/PPTCODES/Source/ExcelFiles/Phase1_SiteVisit_Debugger.xlsx' #Change debugger excel file
wb1 =load_workbook(debugger)
ws1 = wb1.active

#*******Getting all values from excel************
try:
	wb =load_workbook('C:/Users/me/Desktop/PPTCODES/Source/ExcelFiles/Phase1_SiteVisit_Master.xlsx') #change source excel
	ws = wb.active
	print("Excel reading Check....")
	B1314commentcol = 37
	B1314bcommentcol = 39
	B2912commentcol = 38
	for col in range(1,65) :
		
		if(ws.cell(row= 1,column = col).value == 'Region'):
			regioncol = col
		elif(ws.cell(row= 1,column = col).value == 'ED_MARKET'):
			edmarketcol = col
		elif(ws.cell(row= 1,column = col).value == 'SUBMARKET'):
			submarketcol = col
		elif(ws.cell(row= 1,column = col).value == 'County'):
			countycol = col
		elif(ws.cell(row= 1,column = col).value == 'FA_LOCATION'):
			faloccol = col
		elif(ws.cell(row= 1,column = col).value == 'FA.Sector'):
			faseccol = col
		elif(ws.cell(row= 1,column = col).value == 'No.Sectors'):
			no.seccol = col
		elif(ws.cell(row= 1,column = col).value == 'CELL'):
			cellcol = col
		elif(ws.cell(row= 1,column = col).value == 'Band'):
			bandcol = col
		elif(ws.cell(row= 1,column = col).value == 'Structure Type'):
			structurecol = col
		elif(ws.cell(row= 1,column = col).value == 'Total Pictures'):
			totpicscol = col
		elif(ws.cell(row= 1,column = col).value == 'Total Pictures Utilized'):
			utpicscol = col
		elif(ws.cell(row= 1,column = col).value == 'Proxy PIM'):
			proxypimcol = col
		elif(ws.cell(row= 1,column = col).value == 'Capacity'):
			capacitycol = col
		elif(ws.cell(row= 1,column = col).value == 'Noise Rise'):
			noiserisecol = col
		elif(ws.cell(row= 1,column = col).value == 'PIM Potential'):
			pimpotentialcol = col
		elif(ws.cell(row= 1,column = col).value == '95th Percentile DL Throughput'):
			ninteyfifthdlthptcol = col
		elif(ws.cell(row= 1,column = col).value == '95th Percentile RRC Users'):
			ninteyfifthrrccol = col	
		elif(ws.cell(row= 1,column = col).value == '95th Percentile PRB Utilization'):
			ninteyfifthprbcol = col
		elif(ws.cell(row= 1,column = col).value == 'Thpt 95th Percent Fixed'):
			ninteyfifthfixedcol = col	
		elif(ws.cell(row= 1,column = col).value == 'Thpt 95th Percent Not Fixed'):
			ninteyfifthnotfixedcol = col	
		elif(ws.cell(row= 1,column = col).value == 'Sector comments'):
			sectorcommentcol = col	
		elif(ws.cell(row= 1,column = col).value == 'Pim category'):
			pimcatcol = col
		elif(ws.cell(row= 1,column = col).value == 'PIM Potential-Site Index (dB)'):
			pimsiteindexcol = col
		#****************try to Optimize it (Hard Coded) *************
		elif(ws.cell(row= 1,column = col).value == 'C1'):
			C1col = col
		elif(ws.cell(row= 1,column = col).value == 'C2'):
			C2col = col
		elif(ws.cell(row= 1,column = col).value == 'C3'):
			C3col = col
		elif(ws.cell(row= 1,column = col).value == 'C4'):
			C4col = col	
		elif(ws.cell(row= 1,column = col).value == 'C5'):
			C5col = col
		elif(ws.cell(row= 1,column = col).value == 'C6'):
			C6col = col	
		elif(ws.cell(row= 1,column = col).value == 'C7'):
			C7col = col	
		elif(ws.cell(row= 1,column = col).value == 'C8'):
			C8col = col
		elif(ws.cell(row= 1,column = col).value == 'C9'):
			C9col = col
		elif(ws.cell(row= 1,column = col).value == 'C10'):
			C10col = col
		elif(ws.cell(row= 1,column = col).value == 'C11'):
			C11col = col	
		elif(ws.cell(row= 1,column = col).value == 'C12'):
			C12col = col
		elif(ws.cell(row= 1,column = col).value == 'B13/14'):
			B1314commentcol = col
		elif(ws.cell(row= 1,column = col).value == 'B13/14.b'):
			B1314bcommentcol = col
		elif(ws.cell(row= 1,column = col).value == 'B29/12'):
			b2912commentcol = col
		
		com = []
		sitecomments = []
		

	#************************************************
	
	#Taking template.pptx and start working on it

	prs = Presentation('C:/Users/me/Desktop/PPTCODES/Source/Template.pptx')
	#Editing slide1 from the template read
	Slide1 = prs.slides[0]
	pic = Slide1.shapes.add_picture(att_path,Inches(12.31),Inches(6.43),Inches(0.92),Inches(0.92))
	pic = Slide1.shapes.add_picture(tts_path, Inches(0),Inches(0),Inches(13.33),Inches(0.92))
	Slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.38), Inches(1.86), Inches(11.85), Inches(2.02))
	Slide1.shapes[2].text_frame.paragraphs[0].text = 'Proactive PIM Assessment for B14'
	Slide1.shapes[2].text_frame.paragraphs[0].font.size = Pt(35)
	Slide1.shapes[2].fill.solid()	
	Slide1.shapes[2].fill.fore_color.rgb = RGBColor(0,199,255)

	Slide1.shapes.add_textbox(Inches(0.38),Inches(4.19),Inches(12),Inches(3))
	Slide1.shapes[3].text_frame.paragraphs[0].text = 'FALOCATION      : ' + sys.argv[1]
	#adding 5 lines to the paragraphs text_frame
	i = 1
	for i in range(1,7):
		Slide1.shapes[3].text_frame.add_paragraph() 
		i = i+1

	#Second Slide 


	Slide2 = prs.slides[1]
	pic = Slide2.shapes.add_picture(att_path,Inches(12.31),Inches(6.43),Inches(0.92),Inches(0.92))
	pic = Slide2.shapes.add_picture(tts_path, Inches(0),Inches(0),Inches(13.33),Inches(0.92))

	Slide2.shapes.add_textbox(Inches(0.15),Inches(0.78),Inches(12.87),Inches(0.74))
	Slide2.shapes[3].text_frame.paragraphs[0].text = sys.argv[2] + '_' + sys.argv[1] + '-' + sys.argv[1]
	Slide2.shapes[3].text_frame.paragraphs[0].font.size = Pt(30)

	Slide2.shapes[3].text_frame.add_paragraph()
	Slide2.shapes[3].text_frame.paragraphs[1].text = 'Summary'
	Slide2.shapes[3].text_frame.paragraphs[1].font.size = Pt(20)
	Slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.09), Inches(1.41), Inches(2.77), Inches(0.32))
	Slide2.shapes[4].text_frame.paragraphs[0].text = 'B17 PROXY PIM'
	Slide2.shapes[4].text_frame.paragraphs[0].font.size = Pt(12)
	Slide2.shapes[4].text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
	Slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.02), Inches(1.41), Inches(2.77), Inches(0.32))
	Slide2.shapes[5].text_frame.paragraphs[0].text = 'PIM POTENTIAL SITE INDEX'
	Slide2.shapes[5].text_frame.paragraphs[0].font.size = Pt(12)
	Slide2.shapes[5].text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
	Slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.94), Inches(1.41), Inches(2.77), Inches(0.32))
	Slide2.shapes[6].text_frame.paragraphs[0].text = 'LOW RISK PIM'
	Slide2.shapes[6].text_frame.paragraphs[0].font.size = Pt(12)
	Slide2.shapes[6].text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)

	#table1 = Slide2.shapes.add_table(5,11,Inches(4.08),Inches(1.78),Inches(8.65),Inches(1.59)).table
	table1 = Slide2.shapes[0].table
	table2 = Slide2.shapes.add_table(int(sectors) + 1,8,Inches(4.09),Inches(3.2),Inches(8.65),Inches(1.62)).table
	table2.cell(0,0).text_frame.paragraphs[0].text = 'Cell'
	table2.cell(0,1).text_frame.paragraphs[0].text = 'Proxy Pim'
	table2.cell(0,2).text_frame.paragraphs[0].text = 'Capacity'
	table2.cell(0,3).text_frame.paragraphs[0].text = 'Noise Rise'
	table2.cell(0,4).text_frame.paragraphs[0].text = 'Proxy Pim Potential'
	table2.cell(0,5).text_frame.paragraphs[0].text = 'RRC 95th Percentile'
	table2.cell(0,6).text_frame.paragraphs[0].text = 'Peak Mbps 95th Percentile'
	table2.cell(0,7).text_frame.paragraphs[0].text = 'PIM Potential Site Index(dB)'

	count = 0
	totpicsbox = 0
	
	for i in range(0,int(sectors) + 1):
		for j in range(0,8):
			str1 = table2.cell(i,j).text_frame.paragraphs[0]
			str1.alignment = PP_ALIGN.CENTER
			_set_cell_border(table2.cell(i,j), border_color="000000", border_width='12700')
			str1.font.size = Pt(11)	
	for m in range(1,int(sectors)+1)	:
		table2.rows[m].height = Inches(0.3)
	#site specific comments
	ListComments = []
	txtbox = Slide2.shapes.add_textbox(Inches(4.09),Inches(4.83),Inches(8.85),Inches(1.5))
	txtbox.line.color.rgb = RGBColor(0,0,0)
	txtbox.text_frame.paragraphs[0].text = 'Site Specific Recommandations'
	txtbox.text_frame.paragraphs[0].font.size = Pt(13)
	txtbox.text_frame.margin_bottom = Inches(0.10)
	txtbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
	txtbox.text_frame.paragraphs[0].font.bold = True
	txtbox.text_frame.paragraphs[0].word_wrap = True
	txtbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
	
	

	#Slide 3 and above for part1
	
	for s in range(1,int(sectors) + 1):
		blank_slide_layout = prs.slide_layouts[6]
		Slide = prs.slides.add_slide(blank_slide_layout)
		pic = Slide.shapes.add_picture(att_path,Inches(12.31),Inches(6.43),Inches(0.92),Inches(0.92))
		pic = Slide.shapes.add_picture(tts_path, Inches(0),Inches(0),Inches(13.33),Inches(0.92))
		Slide.shapes.add_textbox(Inches(0.15),Inches(0.78),Inches(12.87),Inches(0.74))
		Slide.shapes[2].text_frame.paragraphs[0].text =  sys.argv[2] + '_' + sys.argv[1] +  '-' + sys.argv[1] 
		Slide.shapes[2].text_frame.paragraphs[0].font.size = Pt(30)
		Slide.shapes[2].text_frame.add_paragraph()
		Slide.shapes[2].text_frame.paragraphs[1].text = 'INDEX SCORE'
		Slide.shapes[2].text_frame.paragraphs[1].font.size = Pt(20)
		Slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.15), Inches(1.68), Inches(3.94), Inches(0.4))
		Slide.shapes[3].text_frame.paragraphs[0].text = 'B17 PROXY PIM'
		Slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.56), Inches(1.68), Inches(3.94), Inches(0.4))
		Slide.shapes[4].text_frame.paragraphs[0].text = 'PIM POTENTIAL SITE INDEX'
		Slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.97), Inches(1.68), Inches(3.94), Inches(0.4))
		Slide.shapes[5].text_frame.paragraphs[0].text = 'LOW RISK PIM'
		for sh in range(3,6):
			Slide.shapes[sh].text_frame.paragraphs[0].font.size = Pt(12)
			Slide.shapes[sh].text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
		table1 = Slide.shapes.add_table(4,4,Inches(0.15),Inches(2.3),Inches(4.86),Inches(2.16)).table
		#slidey.shapes[6] is kpi table in slides 2-sectors
		table1.cell(0,1).text_frame.paragraphs[0].text = 'Peak RRC 95th Percentile'
		table1.cell(0,2).text_frame.paragraphs[0].text = 'Peak % PRB 95th Percentile'
		table1.cell(0,3).text_frame.paragraphs[0].text = 'Peak Mbps 95th Percentile'
		table1.cell(1,0).text_frame.paragraphs[0].text = 'B17 Cell'
		table1.cell(2,0).text_frame.paragraphs[0].text = 'Potential B14 Performance (if PIM fixed)'
		table1.cell(3,0).text_frame.paragraphs[0].text = 'Potential B14 Performance (if PIM not fixed)'
		for i in range(0,4):
			for j in range(0,4):
				str1 = table1.cell(i,j).text_frame.paragraphs[0]
				str1.alignment = PP_ALIGN.CENTER
				_set_cell_border(table1.cell(i,j), border_color="000000", border_width='12700')
				str1.font.size = Pt(10)
		table2 = Slide.shapes.add_table(2,8,Inches(5.35),Inches(2.29),Inches(7.55),Inches(0.71)).table
		for x in range(0,7):
			for y in range(0,2):
				table2.cell(y,x).text_frame.paragraphs[0].text = 'C' + str(x + 1)
				table2.cell(y,x).text_frame.paragraphs[0].font.bold = False
				table2.cell(y,x).text_frame.paragraphs[0].font.size = Pt(12)
				table2.cell(y,x).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
		table3 = Slide.shapes.add_table(2,5,Inches(5.35),Inches(3.08),Inches(4.66),Inches(0.71)).table
		for x in range(0,5):
			for y in range(0,2):
				table3.cell(y,x).text_frame.paragraphs[0].text = 'C' + str(x + 8)
				table3.cell(y,x).text_frame.paragraphs[0].font.bold = False
				table3.cell(y,x).text_frame.paragraphs[0].font.size = Pt(12)
				table3.cell(y,x).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
		table2.cell(0,7).text_frame.paragraphs[0].text = 'Total'
		table2.cell(0,7).text_frame.paragraphs[0].font.bold = False
		table2.cell(0,7).text_frame.paragraphs[0].font.size = Pt(12)
		table2.cell(0,7).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
		#sector specific recommandation
		txtbox2 = Slide.shapes.add_textbox(Inches(0.15),Inches(4.67),Inches(4.86),Inches(2.76))
		txtbox2.line.color.rgb = RGBColor(0,0,0)
		txtbox2.text_frame.paragraphs[0].text = 'Sector Specific Recommandations'
		txtbox2.text_frame.paragraphs[0].font.size = Pt(15)
		txtbox2.text_frame.word_wrap = True
		txtbox2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
		
		pic = Slide.shapes.add_picture(defpath,Inches(5.27),Inches(4.19),Inches(7.04),Inches(3.09))
		

	#****************** Updating tables with excel values ***********

	sectors = int(sys.argv[3]) + 2
	faloc = sys.argv[1] 
	slide1 = prs.slides[1]
	shape1 = slide1.shapes
	#slide 2 shapes 4,5,6 are tabs
	shape1[4].fill.solid()	
	shape1[4].fill.fore_color.rgb = RGBColor(204,230,255)
	shape1[5].fill.solid()
	shape1[5].fill.fore_color.rgb = RGBColor(204,230,255)
	shape1[6].fill.solid()
	shape1[6].fill.fore_color.rgb = RGBColor(204,230,255)
	
	#-------Variable for the spectrum scan picture comments from excel************
	
	ACommentB1314 = ' '
	ACommentB1314b = ' '
	ACommentB2912 = ' '
	
	BCommentB1314 = ' '
	BCommentB1314b = ' '
	BCommentB2912 = ' '
	
	CCommentB1314 = ' '
	CCommentB1314b = ' '
	CCommentB2912 = ' '
	
	# Slide 3,4,5 for 3 sectors
	
	for y in range(2,sectors):
		slidey = prs.slides[y]
		shapey = slidey.shapes
		shapecount = 0
		shapey[3].fill.solid()	
		shapey[3].fill.fore_color.rgb = RGBColor(204,230,255)
		shapey[4].fill.solid()
		shapey[4].fill.fore_color.rgb = RGBColor(204,230,255)
		shapey[5].fill.solid()
		shapey[5].fill.fore_color.rgb = RGBColor(204,230,255)
				
		x =0
		y = 1
		Acell = str(sys.argv[2])
		Bcell = str(sys.argv[2])
		Ccell = str(sys.argv[2])
		
		
		print(faloc)
		for cell in ws['C'] :
			x = x + 1
			sdicount = 0
			if(str(cell.value) == str(faloc)):
			
				cellValue = str(ws.cell(row= x, column = cellcol).value)
				if cellValue[10] == 'A':
					
					ACommentB1314 = str(ws.cell(row= x, column = 37).value)
					ACommentB1314b = str(ws.cell(row= x, column = 39).value)
					ACommentB2912 = str(ws.cell(row= x, column = 38).value)
				elif cellValue[10] == 'B':
					BCommentB1314 = str(ws.cell(row= x, column = 37).value)
					BCommentB1314b = str(ws.cell(row= x, column = 39).value)
					BCommentB2912 = str(ws.cell(row= x,column = 38).value)
				elif cellValue[10] == 'C':
					CCommentB1314 = str(ws.cell(row= x, column = 37).value)
					CCommentB1314b = str(ws.cell(row= x, column = 39).value)
					CCommentB2912 = str(ws.cell(row= x, column = 38).value)
				else:
					print('cell value not read for spectrum slides')
				totalpics = str(ws.cell(row= x, column = totpicscol).value)
				Picsutilized = str(ws.cell(row= x, column = utpicscol).value)
				#updating 1st slide info
				Slide1.shapes[3].text_frame.paragraphs[1].text = 'Region                 :' + str(ws.cell(row= x, column = regioncol).value)
				Slide1.shapes[3].text_frame.paragraphs[2].text = 'Market                :' + str(ws.cell(row= x, column = edmarketcol).value)
				Slide1.shapes[3].text_frame.paragraphs[3].text = 'Sub Market        :'  + str(ws.cell(row= x, column = submarketcol).value)
				Slide1.shapes[3].text_frame.paragraphs[4].text = 'County                :'    + str(ws.cell(row= x, column = countycol).value)
				#Slide1.shapes[3].text_frame.paragraphs[5].text = 'Structure Type   :'    + str(ws.cell(row= x, column = structurecol).value)
				for para in Slide1.shapes[3].text_frame.paragraphs :
					para.font.size = Pt(25)			
				y = y+1
				slidey = prs.slides[y]
				shapey = slidey.shapes
				Slide2table2 = Slide2.shapes[7].table
				Slide2table2.cell(y - 1,0).text_frame.paragraphs[0].text  = str(ws.cell(row= x, column = cellcol).value)
				Slide2table2.columns[0].width = Inches(1.5)
				
				Slide2table2.cell(y - 1,1).text_frame.paragraphs[0].text  = str(ws.cell(row= x, column = proxypimcol).value)
				Slide2table2.cell(y - 1,2).text_frame.paragraphs[0].text  = str(ws.cell(row= x, column = noiserisecol).value)
				Slide2table2.cell(y - 1,3).text_frame.paragraphs[0].text  = str(ws.cell(row= x, column = capacitycol).value)
				if str(ws.cell(row= x, column = pimpotentialcol).value) == 'None':
					Slide2table2.cell(y - 1,4).text_frame.paragraphs[0].text = ' '
				else:
					Slide2table2.cell(y - 1,4).text_frame.paragraphs[0].text  = str(ws.cell(row= x, column = pimpotentialcol).value)
				Slide2table2.cell(y - 1,5).text_frame.paragraphs[0].text  = str(ws.cell(row= x, column = ninteyfifthrrccol).value)
				Slide2table2.cell(y - 1,6).text_frame.paragraphs[0].text  = str(ws.cell(row= x, column = ninteyfifthdlthptcol).value)
				Slide2table2.cell(y- 1,7).text_frame.paragraphs[0].text = str(ws.cell(row=x,column = pimsiteindexcol).value)
				
				#Updating comments for slide 2 and above
				
				if str(ws.cell(row= x, column = sectorcommentcol).value) == 'None':
					print('sector comment empty')
				else:
					slidey.shapes[9].text_frame.add_paragraph()	
					slidey.shapes[9].text_frame.paragraphs[1].text = str(str(ws.cell(row= x, column = sectorcommentcol).value))
					ListComments.append(str(ws.cell(row= x, column = sectorcommentcol).value)) #for sector comments
				
				com = list(set(ListComments))
				for para in slidey.shapes[9].text_frame.paragraphs :
					para.font.size = Pt(13)
				
				slidey.shapes[2].text_frame.paragraphs[0].text = sys.argv[2] + '_' + sys.argv[1] + '-' + str(ws.cell(row = x, column = cellcol).value)
				
				
				
				if totpicsbox == 0:
					county = Slide2.shapes.add_textbox(Inches(4.09),Inches(6.81),Inches(3.76),Inches(0.46))
					county.text_frame.paragraphs[0].text = 'County :  ' + str(ws.cell(row = x, column = countycol).value)
					county.text_frame.paragraphs[0].font.size = Pt(13)
					county.text_frame.paragraphs[0].font.bold = True
					county.text_frame.add_paragraph()
					county.text_frame.paragraphs[1].text = 'Structure Type :  ' + str(ws.cell(row = x, column = structurecol).value)
					county.text_frame.paragraphs[1].font.size = Pt(13)
					county.text_frame.paragraphs[1].font.bold = True
					county.line.color.rgb = RGBColor(0,0,0)
			
					
					totpicsbox = totpicsbox + 1   #so that this text box will appear only once in slide2
				#updating kpi values in slide2 and above table
				
				kpitable = slidey.shapes[6].table
				kpitable.cell(1,1).text_frame.paragraphs[0].text = str(ws.cell(row = x, column = ninteyfifthrrccol).value)
				kpitable.cell(1,2).text_frame.paragraphs[0].text = str(ws.cell(row = x, column = ninteyfifthprbcol).value)
				kpitable.cell(1,3).text_frame.paragraphs[0].text = str(ws.cell(row = x, column = ninteyfifthdlthptcol).value)
				if str(ws.cell(row = x, column = ninteyfifthfixedcol).value) == 'None':
					kpitable.cell(2,3).text_frame.paragraphs[0].text = ' '
				else:
					kpitable.cell(2,3).text_frame.paragraphs[0].text = str(ws.cell(row = x, column = ninteyfifthfixedcol).value)
				if str(ws.cell(row = x, column = ninteyfifthnotfixedcol).value) == 'None':
					kpitable.cell(3,3).text_frame.paragraphs[0].text = ' '
				else:
					kpitable.cell(3,3).text_frame.paragraphs[0].text = str(ws.cell(row = x, column = ninteyfifthnotfixedcol).value)
				
				
				#********* Category Score ************
				
				pimcat = str(ws.cell(row= x, column = pimcatcol).value)
				print(str(pimcat))
				cat1 = str(ws.cell(row= x,column = C1col).value)
				cat2 = str(ws.cell(row= x,column = C2col).value)
				cat3 = str(ws.cell(row= x,column = C3col).value)
				cat4 = str(ws.cell(row= x,column = C4col).value)
				cat5 = str(ws.cell(row= x,column = C5col).value)
				cat6 = str(ws.cell(row= x,column = C6col).value)
				cat7 = str(ws.cell(row= x,column = C7col).value)
				cat8 = str(ws.cell(row= x,column = C8col).value)
				cat9 = str(ws.cell(row= x,column = C9col).value)
				cat10 = str(ws.cell(row= x,column = C10col).value)
				cat11 = str(ws.cell(row= x,column = C11col).value)
				cat12 = str(ws.cell(row= x,column = C12col).value)
				
				
				sdi = int(cat1) + int(cat2) + int(cat3) + int(cat4) + int(cat5) + int(cat6) + int(cat7) + int(cat8) + int(cat9) + int(cat10) + int(cat11) + int(cat12)
				print(str(sdi))
				
				#*************************************
				# table1 = slide2.shapes[7].table # kpi table in every sector slide
				scoretable1 = shapey[7].table
				scoretable2 = shapey[8].table
				
				scoretable1.cell(1,0).text_frame.paragraphs[0].text = cat1
				scoretable1.cell(1,1).text_frame.paragraphs[0].text = cat2
				scoretable1.cell(1,2).text_frame.paragraphs[0].text = cat3
				scoretable1.cell(1,3).text_frame.paragraphs[0].text = cat4
				scoretable1.cell(1,4).text_frame.paragraphs[0].text = cat5
				scoretable1.cell(1,5).text_frame.paragraphs[0].text = cat6
				scoretable1.cell(1,6).text_frame.paragraphs[0].text = cat7
				
				scoretable2.cell(1,0).text_frame.paragraphs[0].text = cat8
				scoretable2.cell(1,1).text_frame.paragraphs[0].text = cat9
				scoretable2.cell(1,2).text_frame.paragraphs[0].text = cat10
				scoretable2.cell(1,3).text_frame.paragraphs[0].text = cat11
				scoretable2.cell(1,4).text_frame.paragraphs[0].text = cat12
				
				if sdicount == 0:
					sditextbox = slidey.shapes.add_textbox(Inches(5.27),Inches(3.87),Inches(1.98),Inches(0.38))
					sditextbox.text_frame.paragraphs[0].text = 'Site Data Index : ' + str(sdi)
					sditextbox.text_frame.paragraphs[0].font.size = Pt(15)
					scoretable1.cell(1,7).text_frame.paragraphs[0].text = str(sdi)
					sdicount = sdicount + 1
					
					
				
				for b in range(0,8):
					s1 = scoretable1.cell(1,b).text_frame.paragraphs[0]
					s1.alignment = PP_ALIGN.CENTER
					s1.font.size = Pt(10)
				
				for d in range(0,5):
					s2 = scoretable2.cell(1,d).text_frame.paragraphs[0]
					s1.alignment = PP_ALIGN.CENTER
					s2.font.size = Pt(10)
				
				
				#*************************************
				#print('pim cat:' + pimcat)
				shapey[3].fill.solid()
				shapey[4].fill.solid()
				shapey[5].fill.solid()
				if(pimcat == 'B17 PROXY PIM') :
					shape1[4].fill.fore_color.rgb = RGBColor(77,166,255)
					shapey[3].fill.fore_color.rgb = RGBColor(77,166,255)
					shapey[4].fill.fore_color.rgb = RGBColor(204,230,255)
					print(pimcat)			
				elif(pimcat == 'POTENTIAL PIM SITE INDEX') :	
					shapey[4].fill.fore_color.rgb = RGBColor(77,166,255)
					shape1[5].fill.fore_color.rgb = RGBColor(77,166,255)
					print(pimcat)
				elif(pimcat == 'LOW RISK PIM ') :
					
					shapey[5].fill.fore_color.rgb = RGBColor(77,166,255)
					shape1[6].fill.fore_color.rgb = RGBColor(77,166,255)
					print(pimcat)		
				else :
					print('Tabs not detected ')

	#combining sector comments to site comments
	com = set(com)
	print(com)
	print('******************')
	#scomments = str(com)
	str(com).splitlines()
	print(com)
	scomments = '\n'.join(com)
	#print(scomments)
	finalcomment = str(scomments)
	#print(finalcoment)
	k = finalcomment.splitlines()
	#print(k)
	k = set(k)
	sk = '\n'.join(k)
	sk = sk.splitlines()
	print('Comments length' + str(len(sk)))
	#print(sk)
	print('**********************************************************')
	#print(sk)
	#print(str(scomments))
	#print(Slide2.shapes[8].text_frame.paragraphs[0].text )
	for i in range(0,len(sk)):
		Slide2.shapes[8].text_frame.add_paragraph()		
		Slide2.shapes[8].text_frame.paragraphs[i+1].text  = sk[i]
		Slide2.shapes[8].text_frame.paragraphs[i+1].font.size = Pt(13)
	#*************************************** PART 2 ******************************************8

	if os.path.exists(dir_path):
	# ........ Getting Values for Second Half .............

		countA = 0
		countB = 0
		countG = 0
		countM = 0
		countS = 0
		ListA = []
		ListB = []
		ListG = []
		ListM = []
		ListS = []
		Total_list = []
	 #print('Created lists')
	# Reading each files in the Source directory and splitting them with "_"



		for file in os.listdir(dir_path) :
			print(file)
			temp = file.split("_")
	#Assinging all the photos of Alphas and betas and gammas respective variables
		 
	#***** Checking for Alpha pics ***********
			if temp[1].upper() == 'A.JPG' or temp[1].upper() == 'A1.JPG' or temp[1].upper() == 'A.PNG' or temp[1].upper() == 'A1.PNG' or temp[1].upper() == 'A1' or temp[1].upper() == 'A2' or temp[1].upper() == 'A3' or temp[1].upper() == 'A4' or temp[1].upper() == 'A5' or temp[1].upper() == 'A6' or temp[1].upper() == 'A7' or temp[1].upper() == 'A8' or temp[1].upper() == 'A9' :
				countA = countA+1
				ListA.append(file)
	   
	# ******* Beta Checking *************
			if temp[1].upper() == 'B.JPG' or temp[1].upper() == 'B1.JPG' or temp[1].upper() == 'b.PNG' or temp[1].upper() == 'B1.PNG' or temp[1].upper() == 'B1' or temp[1].upper() == 'B2' or temp[1].upper() == 'B3' or temp[1].upper() == 'B4' or temp[1].upper() == 'B5' or temp[1].upper() == 'B6' or temp[1].upper() == 'B7' or temp[1].upper() == 'B8' or temp[1].upper() == 'B9' :
				countB = countB+1
				ListB.append(file)

	#******* Checking for Gamma *********
			if temp[1].upper() == 'G.JPG' or temp[1].upper() == 'G1.JPG' or temp[1].upper() == 'G.PNG' or temp[1].upper() == 'G1.PNG' or temp[1].upper() == 'G1' or temp[1].upper() == 'G2' or temp[1].upper() == 'G3' or temp[1].upper() == 'G4' or temp[1].upper() == 'G5' or temp[1].upper() == 'G6' or temp[1].upper() == 'G7' or temp[1].upper() == 'G8' or temp[1].upper() == 'G9':
				countG = countG+1
				ListG.append(file)
	   

	  
	# End of Checking for sector pics

	#Checking for rfds and siteplans & sitepics	
		if os.path.exists(rfds_path):
			for file2 in os.listdir(rfds_path) :
				print(file2)
				temp = file2.split("_")
				#****** Check for RFDS ******
				if temp[1].upper() == 'RFDS.PNG' or temp[1].upper() == 'RFDS.JPG':
					countM = countM+1
					ListM.append(file2)
		   
				try:
					if temp[1].upper() == 'A' and (temp[2].upper() == 'RFDS.PNG' or temp[2].upper() == 'RFDS.JPG'):
						countM = countM+1
						ListM.append(file2)
					if temp[1].upper() == 'B' and (temp[2].upper() == 'RFDS.PNG' or temp[2].upper() == 'RFDS.JPG'):
						countM = countM+1
						ListM.append(file2)
					if temp[1].upper() == 'G' and (temp[2].upper() == 'RFDS.PNG' or temp[2].upper() == 'RFDS.JPG'):
						countM = countM+1
						ListM.append(file2)
				except:
					pass
		   
		   
		#***** Checking for Siteplan and SitePic *****
				if temp[1].upper() == 'SITEPLAN.PNG' or temp[1].upper() == 'SITEPLAN.JPG':
					countS = countS+1
					ListS.append(file2)
				if temp[1].upper() == 'SITEPIC.JPG' or temp[1].upper() == 'SITEPIC.PNG':
					countS = countS+1
					ListS.append(file2)
		 
		else :
			print('RFDS and SITEPLANS NOT AVIALABLE')	


	#Total_list.append(temp[1])
		print('count of Alpha pics :' + str(countA))
		print('count of Beta pics :'  + str(countB))
		print('count of Gamma pics :' + str(countG))
		print(ListA)
		print(ListB)
		print(ListG)
		print(ListM)
		print(ListS)
	 
	 
	 

	# ............ Core ......................

		countAx =0
		countBx =0
		countGx =0
		countMx =0
	 
	#***** Site Plan Sitepics Slide *********
		s = 0
		while s < len(ListS):
			dir_pathS1 = rfds_path + '/' + ListS[s]

			print('Created SITEPLAN Slide')
			
			blank_slide_layout = prs.slide_layouts[6]
			Slide = prs.slides.add_slide(blank_slide_layout)
			pic = Slide.shapes.add_picture(att_path,Inches(12.31),Inches(6.43),Inches(0.92),Inches(0.92))
			pic = Slide.shapes.add_picture(tts_path, Inches(0),Inches(0),Inches(13.33),Inches(0.92))
	
			txBox = Slide.shapes.add_textbox(Inches(0.46), Inches(0.92), Inches(12.77), Inches(1.05))
			txBox.text_frame.paragraphs[0].text = sys.argv[2] + '_' + sys.argv[1]
			txBox.text_frame.paragraphs[0].font.size = Pt(30)
			txBox.text_frame.add_paragraph()
			txBox.text_frame.paragraphs[1].text = 'SITE PLAN'
			txBox.text_frame.paragraphs[1].font.size = Pt(20)
			tablepics = Slide.shapes.add_table(2,2,Inches(6.76),Inches(0.82),Inches(5.87),Inches(0.6)).table
			tablepics.cell(0,0).text_frame.paragraphs[0].text = 'Total Pictures Available'
			tablepics.cell(0,0).text_frame.paragraphs[0].font.size = Pt(14)
			tablepics.cell(1,0).text_frame.paragraphs[0].text = 'Total Pictures Utilized'
			tablepics.cell(1,0).text_frame.paragraphs[0].font.size = Pt(14)
			tablepics.cell(0,1).text_frame.paragraphs[0].text = totalpics
			tablepics.cell(1,1).text_frame.paragraphs[0].text = Picsutilized
			tablepics.cell(0,1).text_frame.paragraphs[0].font.size = Pt(14)
			tablepics.cell(1,1).text_frame.paragraphs[0].font.size = Pt(14)
			tablepics.cell(0,1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
			tablepics.cell(1,1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
			#inserting siteplan pics
			try:
				pic = Slide.shapes.add_picture(dir_pathS1, left=left, top=top, width= width,  height=height)
				if(s+1 >= len(ListS)):
					break
				else:
					dir_pathS2 = rfds_path + '/' + ListS[s+1]
					pic = Slide.shapes.add_picture(dir_pathS2, left=Inches(6.67), top=top, width= width, height=height)
					s = s+2
			except Exception as e:
				s = s+2
				print('problem with SITEPICS AND SIETPLANS')
				ws1.cell(row = ws1.max_row + 1 , column = 1).value = sys.argv[1]
				ws1.cell(row = ws1.max_row  , column = 2).value = str('Insert SITEPICS AND SITEPLAN')
				ws1.cell(row = ws1.max_row  , column = 3).value = str(e)
				pass
				
	 
	#***** Alpha Sector (part2) *********
		n = 0
		while n < len(ListA):
			dir_pathA1 = dir_path + '/' + ListA[n]
			print('Created Alpha new Slide')
			blank_slide_layout = prs.slide_layouts[6]
			Slide = prs.slides.add_slide(blank_slide_layout)
			pic = Slide.shapes.add_picture(att_path,Inches(12.31),Inches(6.43),Inches(0.92),Inches(0.92))
			pic = Slide.shapes.add_picture(tts_path, Inches(0),Inches(0),Inches(13.33),Inches(0.92))

			txBox = Slide.shapes.add_textbox(Inches(0.46), Inches(0.92), Inches(12.77), Inches(1.05))
			tf = txBox.text_frame
			tf.word_wrap = True
			tf.paragraphs[0].text = sys.argv[2] + '_' + sys.argv[1] + '-' + sys.argv[1] + '.A'
			tf.paragraphs[0].font.size = Pt(30)
			tf.add_paragraph()
			tf.paragraphs[1].text = 'Site Data Analytics: Pictures'
			tf.paragraphs[1].font.size = Pt(20)
	# images from alpha goes here!
			try:
				pic = Slide.shapes.add_picture(dir_pathA1, left=left, top=top, width= width,  height=height)
				GetPicComment(pic._pic.nvPicPr.cNvPr.get('descr'),'1')
				
			except Exception as e:
				
				n = n +1
				print('problem with alpha one')
				ws1.cell(row = ws1.max_row + 1 , column = 1).value = sys.argv[1]
				ws1.cell(row = ws1.max_row  , column = 2).value = str('Insert picture Alpha: at 1st position')
				ws1.cell(row = ws1.max_row  , column = 3).value = str(e)
				
				pass
	  
			if(n+1 >= len(ListA)):
				break
			else:
				dir_pathA2 = dir_path + '/' + ListA[n+1]
				try:
					pic = Slide.shapes.add_picture(dir_pathA2, left=Inches(6.67), top=top, width= width, height=height)
					n = n+2
					GetPicComment(pic._pic.nvPicPr.cNvPr.get('descr'),'2')
					#//////////////////////////////////////////////////////////////////////
				except Exception as e:
					n= n+2
					print('problem with alpha 2')
					ws1.cell(row = ws1.max_row + 1 , column = 1).value = sys.argv[1]
					ws1.cell(row = ws1.max_row  , column = 2).value = str('Insert picture Alpha: at 2nd position')
					ws1.cell(row = ws1.max_row  , column = 3).value = str(e)
					pass

	#***** Beta Sector(Part 2) *********
		m = 0
		while m < len(ListB):
			dir_pathB1 = dir_path + '/' + ListB[m]
			print('Created Beta new Slide')
			blank_slide_layout = prs.slide_layouts[6]
			Slide = prs.slides.add_slide(blank_slide_layout)
			pic = Slide.shapes.add_picture(att_path,Inches(12.31),Inches(6.43),Inches(0.92),Inches(0.92))
			pic = Slide.shapes.add_picture(tts_path, Inches(0),Inches(0),Inches(13.33),Inches(0.92))
			txBox = Slide.shapes.add_textbox(Inches(0.46), Inches(0.92), Inches(12.77), Inches(1.05))
			tf = txBox.text_frame
			tf.word_wrap = True
			tf.paragraphs[0].text = sys.argv[2] + '_' + sys.argv[1] + '-' + sys.argv[1] + '.B'
			tf.paragraphs[0].font.size = Pt(30)
			tf.add_paragraph()
			tf.paragraphs[1].text = 'Site Data Analytics: Pictures'
			tf.paragraphs[1].font.size = Pt(20)

	#Inserting pics for Beta Sector!!
			try:
				pic = Slide.shapes.add_picture(dir_pathB1, left=left, top=top, width= width,  height=height)
				GetPicComment(pic._pic.nvPicPr.cNvPr.get('descr'),'1')
			except Exception as e:
				m = m +1
				print('problem with Beta one')
				ws1.cell(row = ws1.max_row + 1 , column = 1).value = sys.argv[1]
				ws1.cell(row = ws1.max_row  , column = 2).value = str('Insert picture Beta: at 1st position')
				ws1.cell(row = ws1.max_row  , column = 3).value = str(e)
				pass
				
			if(m+1 >= len(ListB)):
				break
			else:
				dir_pathB2 = dir_path + '/' + ListB[m+1]
				try:
					pic = Slide.shapes.add_picture(dir_pathB2, left=Inches(6.67), top=top, width= width, height=height)
					m = m+2
					GetPicComment(pic._pic.nvPicPr.cNvPr.get('descr'),'2')
					
				except Exception as e: 
					print('problem with beata 2')
					ws1.cell(row = ws1.max_row + 1 , column = 1).value = sys.argv[1]
					ws1.cell(row = ws1.max_row  , column = 2).value = str('Insert picture Beta: at 2nd position')
					ws1.cell(row = ws1.max_row  , column = 3).value = str(e)
					m = m+ 2
					pass
	  
		  
	#*******Gamma Sector(Part 2)*********
		q = 0
		while q < len(ListG):
			dir_pathG1 = dir_path + '/' + ListG[q]

			print('Created Gamma new Slide')
			blank_slide_layout = prs.slide_layouts[6]
			Slide = prs.slides.add_slide(blank_slide_layout)
			pic = Slide.shapes.add_picture(att_path,Inches(12.31),Inches(6.43),Inches(0.92),Inches(0.92))
			pic = Slide.shapes.add_picture(tts_path, Inches(0),Inches(0),Inches(13.33),Inches(0.92))
			txBox = Slide.shapes.add_textbox(Inches(0.46), Inches(0.92), Inches(12.77), Inches(1.05))
			tf = txBox.text_frame
			tf.word_wrap = True
			tf.paragraphs[0].text = sys.argv[2] + '_' + sys.argv[1] + '-' + sys.argv[1] + '.C'
			tf.paragraphs[0].font.size = Pt(30)
			tf.add_paragraph()
			tf.paragraphs[1].text = 'Site Data Analytics: Pictures'
			tf.paragraphs[1].font.size = Pt(20)
			#Inserting pics for Gamma Sector 
			
			try:
				pic = Slide.shapes.add_picture(dir_pathG1, left=left, top=top, width= width,  height=height)
				print(pic._pic.nvPicPr.cNvPr.get('descr'))
				GetPicComment(pic._pic.nvPicPr.cNvPr.get('descr'),'1')
			except Exception as e:
				traceback.print_exc()
				q = q+1
				print('problem with gammaa one')
				ws1.cell(row = ws1.max_row + 1 , column = 1).value = sys.argv[1]
				ws1.cell(row = ws1.max_row  , column = 2).value = str('Insert picture Gamma: at this 1st position')
				ws1.cell(row = ws1.max_row  , column = 3).value = str(e)
				pass
				
			if(q+1 >= len(ListG)):
				break
			else:
				dir_pathG2 = dir_path + '/' + ListG[q+1]
				try:
					pic = Slide.shapes.add_picture(dir_pathG2, left=Inches(6.67), top=top, width= width, height=height)
					q = q+2
					GetPicComment(str(pic._pic.nvPicPr.cNvPr.get('descr')),'2')
				except Exception as e:
					q = q+2
					print('problem with Gamma 2')
					ws1.cell(row = ws1.max_row + 1 , column = 1).value = sys.argv[1]
					ws1.cell(row = ws1.max_row  , column = 2).value = str('Insert picture Gamma at 2nd position')
					ws1.cell(row = ws1.max_row  , column = 3).value = str(e)
					pass

	#***** Pictures for RFDS  *********
	 
		r = 0
		while r < len(ListM):
			dir_pathM1 = rfds_path + '/' + ListM[r]

			print('Created RFDS Slide')

			blank_slide_layout = prs.slide_layouts[6]
			Slide = prs.slides.add_slide(blank_slide_layout)
	 # att and tts template icons
			pic = Slide.shapes.add_picture(att_path,Inches(12.31),Inches(6.43),Inches(0.92),Inches(0.92))
			pic = Slide.shapes.add_picture(tts_path, Inches(0),Inches(0),Inches(13.33),Inches(0.92))
	  # Adding title with text box
			txBox = Slide.shapes.add_textbox(Inches(0.46), Inches(0.92), Inches(12.77), Inches(1.05))
			txBox.text_frame.paragraphs[0].text = sys.argv[1] + '_' + sys.argv[2]
			txBox.text_frame.paragraphs[0].font.size = Pt(30)
			txBox.text_frame.add_paragraph()
			txBox.text_frame.paragraphs[1].text = 'RFDS'
			txBox.text_frame.paragraphs[1].font.size = Pt(20)
	  
	  # images from RFDS goes here!
			pic = Slide.shapes.add_picture(dir_pathM1, left=left, top=top, width= width,  height=height)
			if(r+1 >= len(ListM)):
				break
			else:
				dir_pathM2 = rfds_path + '/' + ListM[r+1]
				pic = Slide.shapes.add_picture(dir_pathM2, left=Inches(6.67), top=top, width= width, height=height)
				r = r+2


		
	# Almost Done
	
		

	else :
		print(sys.argv[1] + " Images for this FALOC doesnot exist in the source " + market + " Market")
		print('Created Thank you slide')

		

	#**********************
	#Ranking
	slide2 = prs.slides[1]  
	shapes = slide2.shapes
	table1 = slide2.shapes[0].table
	table2 = slide2.shapes[7].table



	for a in range(1,sectors - 1):
		for b in range(1,4):
			if(table2.cell(a,b).text_frame.paragraphs[0].text == 'Low'):
				table2.cell(a,b).fill.solid()
				table2.cell(a,b).fill.fore_color.rgb = RGBColor(169,209,142)
			elif(table2.cell(a,b).text_frame.paragraphs[0].text == 'Medium'):
				table2.cell(a,b).fill.solid()
				table2.cell(a,b).fill.fore_color.rgb = RGBColor(237,125,49)
			elif(table2.cell(a,b).text_frame.paragraphs[0].text == 'High'):
				table2.cell(a,b).fill.solid()
				table2.cell(a,b).fill.fore_color.rgb = RGBColor(220,20,60)
				
				

	#******** conditions for all combinitions in table 2  ******

	for x in range(1,sectors - 1):
		ppim1 = table2.cell(x,1).text_frame.paragraphs[0].text 
		cap1  = table2.cell(x,2).text_frame.paragraphs[0].text
		nr1   = table2.cell(x,3).text_frame.paragraphs[0].text



		if ppim1 == 'High' and cap1 == 'High' and nr1 == 'High' :
			val1 = int(table1.cell(2,2).text_frame.paragraphs[0].text) + 1
			table1.cell(2,2).text_frame.paragraphs[0].text = str(val1)
			#print(str(val1))
		elif ppim1 == 'High' and cap1 == 'Medium' and nr1 == 'High' :
			val2 = int(table1.cell(2,3).text_frame.paragraphs[0].text) + 1
			table1.cell(2,3).text_frame.paragraphs[0].text = str(val2)
			#print(str(val2))
		elif ppim1 == 'High' and cap1 == 'Low' and nr1 == 'High' :
			val3 = int(table1.cell(2,4).text_frame.paragraphs[0].text) + 1
			table1.cell(2,4).text_frame.paragraphs[0].text = str(val3)
			#print(str(val3))
		elif ppim1 == 'High' and cap1 == 'High' and nr1 == 'Medium' :
			val4 = int(table1.cell(3,2).text_frame.paragraphs[0].text) + 1
			table1.cell(3,2).text_frame.paragraphs[0].text = str(val4)
			#print(str(val4))
		elif ppim1 == 'High' and cap1 == 'Medium' and nr1 == 'Medium' :
			val5 = int(table1.cell(3,3).text_frame.paragraphs[0].text) + 1
			table1.cell(3,3).text_frame.paragraphs[0].text = str(val5)
			#print(str(val5))
		elif ppim1 == 'High' and cap1 == 'Low' and nr1 == 'Medium' :
			val6 = int(table1.cell(3,4).text_frame.paragraphs[0].text) + 1
			table1.cell(3,4).text_frame.paragraphs[0].text = str(val6)
			#print(str(val6))
		elif ppim1 == 'High' and cap1 == 'High' and nr1 == 'Low' :
			val7 = int(table1.cell(4,2).text_frame.paragraphs[0].text) + 1
			table1.cell(4,2).text_frame.paragraphs[0].text = str(val7)
			#print(str(val7))
		elif ppim1 == 'High' and cap1 == 'Medium' and nr1 == 'Low' :
			val8 = int(table1.cell(4,3).text_frame.paragraphs[0].text) + 1
			table1.cell(4,3).text_frame.paragraphs[0].text = str(val8)
			#print(str(val7))
		elif ppim1 == 'High' and cap1 == 'Low' and nr1 == 'Low' :
			val9 = int(table1.cell(4,4).text_frame.paragraphs[0].text) + 1
			table1.cell(4,4).text_frame.paragraphs[0].text = str(val9)
			#print(str(val9))


		# ******** second set ******************
		if ppim1 == 'Medium' and cap1 == 'High' and nr1 == 'High' :
			val10 = int(table1.cell(2,5).text_frame.paragraphs[0].text) + 1
			table1.cell(2,5).text_frame.paragraphs[0].text = str(val10)
			#print(str(val10))
		elif ppim1 == 'Medium' and cap1 == 'Medium' and nr1 == 'High' :
			val11 = int(table1.cell(2,6).text_frame.paragraphs[0].text) + 1
			table1.cell(2,6).text_frame.paragraphs[0].text = str(val11)
			#print(str(val11))
		elif ppim1 == 'Medium' and cap1 == 'Low' and nr1 == 'High' :
			val12 = int(table1.cell(2,7).text_frame.paragraphs[0].text) + 1
			table1.cell(2,7).text_frame.paragraphs[0].text = str(val12)
			#print(str(val12))
		elif ppim1 == 'Medium' and cap1 == 'High' and nr1 == 'Medium' :
			val13 = int(table1.cell(3,5).text_frame.paragraphs[0].text) + 1
			table1.cell(3,5).text_frame.paragraphs[0].text = str(val13)
			#print(str(val13))
		elif ppim1 == 'Medium' and cap1 == 'Medium' and nr1 == 'Medium' :
			val14 = int(table1.cell(3,6).text_frame.paragraphs[0].text) + 1
			table1.cell(3,6).text_frame.paragraphs[0].text = str(val14)
			#print(str(val14))
		elif ppim1 == 'Medium' and cap1 == 'Low' and nr1 == 'Medium' :
			val15 = int(table1.cell(3,7).text_frame.paragraphs[0].text) + 1
			table1.cell(3,7).text_frame.paragraphs[0].text = str(val15)
			#print(str(val15))
		elif ppim1 == 'Medium' and cap1 == 'High' and nr1 == 'Low' :
			val16 = int(table1.cell(4,5).text_frame.paragraphs[0].text) + 1
			table1.cell(4,5).text_frame.paragraphs[0].text = str(val16)
			#print(str(val16))
		elif ppim1 == 'Medium' and cap1 == 'Medium' and nr1 == 'Low' :
			val17 = int(table1.cell(4,6).text_frame.paragraphs[0].text) + 1
			table1.cell(4,6).text_frame.paragraphs[0].text = str(val17)
			#print(str(val17))
		elif ppim1 == 'Medium' and cap1 == 'Low' and nr1 == 'Low' :
			val18 = table1.cell(4,7).text_frame.paragraphs[0].text
			if val18 == '0':
				val18 = '1'
			elif val18 == '1':
				val18 = '2'
			elif val18 == '2':
				val18 = '3'
			
		

		# ************* third set **************

		if ppim1 == 'Low' and cap1 == 'High' and nr1 == 'High' :
			val19 = int(table1.cell(2,8).text_frame.paragraphs[0].text) + 1
			table1.cell(2,8).text_frame.paragraphs[0].text = str(val19)
			#print(str(val19))
		elif ppim1 == 'Low' and cap1 == 'Medium' and nr1 == 'High' :
			val20 = int(table1.cell(2,9).text_frame.paragraphs[0].text) + 1
			table1.cell(2,9).text_frame.paragraphs[0].text = str(val20)
			#print(str(val20))
		elif ppim1 == 'Low' and cap1 == 'Low' and nr1 == 'High' :
			val21 = int(table1.cell(2,10).text_frame.paragraphs[0].text) + 1
			table1.cell(2,10).text_frame.paragraphs[0].text = str(val21)
			#print(str(val21))
		elif ppim1 == 'Low' and cap1 == 'High' and nr1 == 'Medium' :
			val22 = int(table1.cell(3,8).text_frame.paragraphs[0].text) + 1
			table1.cell(3,8).text_frame.paragraphs[0].text = str(val22)
			#print(str(val22))
		elif ppim1 == 'Low' and cap1 == 'Medium' and nr1 == 'Medium' :
			val23 = int(table1.cell(3,9).text_frame.paragraphs[0].text) + 1
			table1.cell(3,9).text_frame.paragraphs[0].text = str(val23)
			#print(str(val23))
		elif ppim1 == 'Low' and cap1 == 'Low' and nr1 == 'Medium' :
			val24 = int(table1.cell(3,10).text_frame.paragraphs[0].text) + 1
			table1.cell(3,10).text_frame.paragraphs[0].text = str(val24)
			#print(str(val24))
		elif ppim1 == 'Low' and cap1 == 'High' and nr1 == 'Low' :
			val25 = int(table1.cell(4,8).text_frame.paragraphs[0].text) + 1
			table1.cell(4,8).text_frame.paragraphs[0].text = str(val25)
			#print(str(val25))
		elif ppim1 == 'Low' and cap1 == 'Medium' and nr1 == 'Low' :
			val26 = int(table1.cell(4,9).text_frame.paragraphs[0].text) + 1
			table1.cell(4,9).text_frame.paragraphs[0].text = str(val26)
			#print(str(val26))
		elif ppim1 == 'Low' and cap1 == 'Low' and nr1 == 'Low' :
			table1.cell(4,10).text_frame.paragraphs[0].text = '1'
				
		
	# ---- slide 2 comments in the table

	count = 0
	for y in range(2,sectors):
		slidey = prs.slides[y]
		shapey = slidey.shapes
		scoretable1 = shapey[7].table
		scoretable2 = shapey[8].table
		table2 = Slide2.shapes[7].table
		count = count + 1
	
	#changing font size to pt(18) for all cells.
			
	for i in range(2,5):
		for j in range(2,11):
			str1 = table1.cell(i,j).text_frame.paragraphs[0]
			str1.font.size = Pt(11)			
	
	
	#adding google earth images
	
	if os.path.exists('C:/Users/me/Desktop/PPTCODES/Source/GOOGLE EARTH PICTURES/'+ sys.argv[1]):
		pic = Slide2.shapes.add_picture(Gepic1_path,Inches(0.26),Inches(1.73),Inches(3.51),Inches(2.55))
		pic = Slide2.shapes.add_picture(Gepic2_path, Inches(0.26),Inches(4.46),Inches(3.51),Inches(2.55)) 
	else:
		ws1.cell(row = ws1.max_row + 1 , column = 1).value = sys.argv[1]
		ws1.cell(row = ws1.max_row  , column = 2).value = 'GoogleEarth Pictures not found'
	#adding slide 2 disclaimer 
	txtboxdesclaimer = Slide2.shapes.add_textbox(Inches(4.09),Inches(6.37),Inches(8.23),Inches(0.4))
	print('textbox on slide 2 added')
	txtboxdesclaimer.text_frame.paragraphs[0].text = 'NOTE: PROXY PIM POTENTIAL IS CALCULATED ON SECTORS WITH “HIGH” PIM; “LOW” AND “MEDIUM” PIM HAVE NO IMPACT'
	txtboxdesclaimer.text_frame.paragraphs[0].font.size = Pt(12)
	SlideCount = getSlidecount()
	print(SlideCount)
	
	
	
	try:
#------------------------------ SITE VISIT PART ------------------------------------------
		#adding default slide for all site visit candidates
		blank_slide_layout = prs.slide_layouts[6]
		Slide = prs.slides.add_slide(blank_slide_layout)
		spectrumchartpath = 'C:/Users/me/Desktop/PPTCODES/spectrumpic.png'
		pic = Slide.shapes.add_picture(spectrumchartpath,0,0,Inches(13.32),Inches(7.5))
		
		#Adding first alpha slide 
		getSiteVisitSlides('A')
		#Adding first Beta slide
		getSiteVisitSlides('B')
		#Adding first Gamma slide
		getSiteVisitSlides('C')
		
		
		print(ACommentB1314)
		print(ACommentB1314b)
		print(ACommentB2912)
		print(BCommentB1314)
		print(BCommentB1314b)
		print(BCommentB2912)
		print(CCommentB1314)
		print(CCommentB1314b)
		print(CCommentB2912)

	except:
		pass
		

#------------------------------------------------------------------------------------------	
	
	#created thank you slide
	print('Created Thank you slide')
	blank_slide_layout = prs.slide_layouts[6]
	Slide = prs.slides.add_slide(blank_slide_layout)
	pic = Slide.shapes.add_picture(thankyoupicpath,0,0,Inches(13.32),Inches(7.5))
	
		
	
	prs.save('C:/Users/me/Desktop/PPTCODES/Source/ATT-PIM-PHASE1-REPORTS-SITEVISIT/' + sys.argv[1] + '_' + sys.argv[2] + '_V2.pptx')
	print("Powerpoint Saved Successfully ")
except Exception as e :
	print('Exception occured')
	ws1.cell(row = ws1.max_row + 1 , column = 1).value = sys.argv[1]
	ws1.cell(row = ws1.max_row  , column = 2).value = str(e)
	traceback.print_exc()
wb1.save(debugger)

		












		
